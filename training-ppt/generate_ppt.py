# -*- coding: utf-8 -*-
"""
Rhautt Nexus / 瑞合数智枢纽 — 软件完整介绍培训 PPT 生成器
基于 python-pptx，视觉遵循 DESIGN.md（瑞合红 #C8102E / 墨黑 #111827 / 克制专业）。
运行：python generate_ppt.py  →  产出 Rhautt-Nexus-培训.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 品牌 tokens
BRAND        = RGBColor(0xC8, 0x10, 0x2E)
BRAND_HOVER  = RGBColor(0xA8, 0x0D, 0x26)
BRAND_SUBTLE = RGBColor(0xFF, 0xF1, 0xF2)
INK          = RGBColor(0x11, 0x18, 0x27)
INK2         = RGBColor(0x37, 0x41, 0x51)
INK3         = RGBColor(0x6B, 0x72, 0x80)
BG           = RGBColor(0xF9, 0xFA, 0xFB)
SURFACE      = RGBColor(0xFF, 0xFF, 0xFF)
LINE         = RGBColor(0xE5, 0xE7, 0xEB)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS      = RGBColor(0x16, 0xA3, 0x4A)
INFO         = RGBColor(0x25, 0x63, 0xEB)
AMBER        = RGBColor(0xD9, 0x77, 0x06)

CN = "PingFang SC"   # 中文主字体（macOS）；Windows 会回退到雅黑
EN = "Inter"

import os
SHOTS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "shots")
PAGENO = [1]   # 封面 = 1；page_header 自动递增

EMU_IN = 914400
SW = 13.333
SH = 7.5

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- 基础工具
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.06):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shp, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, font=CN,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4):
    """runs: str 或 [(txt, {overrides})]"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        parts = item if isinstance(item, list) else [(item, {})]
        for txt, ov in parts:
            r = p.add_run(); r.text = txt
            f = r.font
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.name = ov.get("font", font)
            f.color.rgb = ov.get("color", color)
            # 让中文也套用同一字体
            rPr = r._r.get_or_add_rPr()
            ea = rPr.makeelement(qn('a:ea'), {'typeface': ov.get("font", font)})
            cs = rPr.makeelement(qn('a:cs'), {'typeface': ov.get("font", font)})
            rPr.append(ea); rPr.append(cs)
    return tb


def page_header(s, kicker, title, idx=None):
    if idx is None:
        PAGENO[0] += 1
        idx = PAGENO[0]
    rect(s, 0, 0, SW, 0.14, fill=BRAND)                       # 顶部品牌条
    text(s, 0.9, 0.55, 11, 0.3, kicker.upper(), size=12, color=BRAND, bold=True, font=EN)
    text(s, 0.9, 0.9, 11.5, 0.9, title, size=30, color=INK, bold=True, spacing=1.05)
    rect(s, 0.9, 1.72, 0.62, 0.06, fill=BRAND)                # 标题下红线
    # 页脚
    text(s, 0.9, 7.02, 8, 0.3, "Rhautt Nexus / 瑞合数智枢纽 · Powered by Rysnova",
         size=9, color=INK3)
    text(s, 11.4, 7.02, 1.05, 0.3, f"{idx:02d}", size=9, color=INK3, align=PP_ALIGN.RIGHT)


def browser_frame(s, fx, fy, fw, url):
    bar = 0.36
    rect(s, fx, fy, fw, bar, fill=INK, rounded=True, radius=0.04)
    for i, cc in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(fx+0.18+i*0.2), Inches(fy+0.13), Inches(0.11), Inches(0.11))
        c.fill.solid(); c.fill.fore_color.rgb = cc; c.line.fill.background(); c.shadow.inherit = False
    pill_x = fx + 0.95
    rect(s, pill_x, fy+0.08, fw-1.15, bar-0.16, fill=RGBColor(0x2A,0x31,0x3F), rounded=True, radius=0.3)
    text(s, pill_x+0.18, fy+0.06, fw-1.4, bar-0.12, url, size=9.5, color=RGBColor(0xC8,0xCE,0xD8),
         font=EN, anchor=MSO_ANCHOR.MIDDLE)
    return bar


def anno_list(s, ax, ay, aw, annos, title="界面导览"):
    text(s, ax, ay, aw, 0.35, title, size=13, color=BRAND, bold=True, font=EN)
    ry = ay + 0.53
    for i, (label, desc) in enumerate(annos):
        b = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ax), Inches(ry), Inches(0.34), Inches(0.34))
        b.fill.solid(); b.fill.fore_color.rgb = BRAND; b.line.fill.background(); b.shadow.inherit = False
        text(s, ax, ry-0.01, 0.34, 0.34, str(i+1), size=12, color=WHITE, bold=True,
             font=EN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, ax+0.5, ry-0.03, aw-0.5, 0.35, label, size=13.5, color=INK, bold=True)
        text(s, ax+0.5, ry+0.3, aw-0.5, 0.7, desc, size=11.5, color=INK2, spacing=1.2)
        ry += 0.92


def shot_slide(kicker, title, img_name, subtitle, annos, url, zoom=None):
    """截图页：左侧浏览器框内嵌真实页面截图，右侧界面导览标注。
    zoom = {crop:(l,t,r,b), at:(x,y,w), label:str} 可选局部放大特写。"""
    s = slide()
    page_header(s, kicker, title)
    text(s, 0.9, 1.78, 11.5, 0.4, subtitle, size=13, color=INK2, spacing=1.25)
    fx, fy, fw = 0.9, 2.32, 6.85
    bar = browser_frame(s, fx, fy, fw, url)
    img = os.path.join(SHOTS, img_name)
    pic = s.shapes.add_picture(img, Inches(fx), Inches(fy+bar), width=Inches(fw))
    ih = Emu(pic.height) / EMU_IN
    rect(s, fx, fy+bar, fw, ih, fill=None, line=LINE, line_w=1)
    anno_list(s, 8.55, fy+0.02, 3.85, annos)
    if zoom:
        l, t, r, b = zoom["crop"]
        zx, zy, zw = zoom["at"]
        zp = s.shapes.add_picture(img, Inches(zx), Inches(zy), width=Inches(zw))
        zp.crop_left, zp.crop_top, zp.crop_right, zp.crop_bottom = l, t, r, b
        zh = Emu(zp.height) / EMU_IN
        rect(s, zx, zy, zw, zh, fill=None, line=BRAND, line_w=2)
        if zoom.get("label"):
            text(s, zx, zy - 0.32, zw, 0.3, zoom["label"], size=11, color=BRAND, bold=True)
    return s


def longshot_slide(kicker, title, img_name, subtitle, annos, url):
    """整页长图页：居中展示完整首页滚动全貌 + 右侧分区说明。"""
    s = slide()
    page_header(s, kicker, title)
    text(s, 0.9, 1.78, 11.5, 0.4, subtitle, size=13, color=INK2, spacing=1.25)
    fx, fy = 1.5, 2.35
    bar = 0.3
    fw_bar = 4.4
    rect(s, fx, fy, fw_bar, bar, fill=INK, rounded=True, radius=0.05)
    text(s, fx+0.22, fy+0.02, fw_bar-0.4, bar-0.04, url, size=9, color=RGBColor(0xC8,0xCE,0xD8),
         font=EN, anchor=MSO_ANCHOR.MIDDLE)
    img = os.path.join(SHOTS, img_name)
    pic = s.shapes.add_picture(img, Inches(fx), Inches(fy+bar), height=Inches(4.35))
    pw = Emu(pic.width) / EMU_IN
    rect(s, fx, fy+bar, pw, 4.35, fill=None, line=LINE, line_w=1)
    anno_list(s, 8.55, fy+0.05, 3.85, annos, title="完整首页 · 分区")
    return s


def cards(s, items, top=2.15, cols=2, gap=0.35, left=0.9, card_h=None, bottom=6.7):
    """items: [(标题, 描述[, 角标])] 网格卡片"""
    n = len(items)
    rows = (n + cols - 1) // cols
    total_w = SW - left * 2
    cw = (total_w - gap * (cols - 1)) / cols
    avail = bottom - top
    ch = card_h if card_h else (avail - gap * (rows - 1)) / rows
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = left + c * (cw + gap)
        y = top + r * (ch + gap)
        card = rect(s, x, y, cw, ch, fill=SURFACE, line=LINE, line_w=1, rounded=True, radius=0.05)
        rect(s, x, y, 0.07, ch, fill=BRAND, rounded=False)     # 左侧品牌竖条
        tag = it[2] if len(it) > 2 else None
        ty = y + 0.22
        if tag:
            text(s, x + 0.32, ty, 1.2, 0.3, tag, size=11, color=BRAND, bold=True, font=EN)
            ty += 0.34
        text(s, x + 0.32, ty, cw - 0.55, 0.5, it[0], size=16, color=INK, bold=True)
        text(s, x + 0.32, ty + 0.42, cw - 0.55, ch - (ty - y) - 0.5, it[1],
             size=12.5, color=INK2, spacing=1.25)
    return cw, ch


# ================================================================ 幻灯片
# ---- 01 封面
s = slide(INK)
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, 0.22, SH, fill=BRAND)
text(s, 1.1, 2.15, 11, 0.4, "企业软件平台 · 完整介绍培训", size=15,
     color=RGBColor(0xC8, 0xCE, 0xD8), bold=True, font=EN)
text(s, 1.05, 2.75, 11.2, 1.6, [
    [("Rhautt Nexus", {"color": WHITE, "size": 54, "bold": True, "font": EN})],
    [("瑞合数智枢纽", {"color": WHITE, "size": 40, "bold": True})],
], spacing=1.05)
rect(s, 1.1, 4.9, 0.9, 0.07, fill=BRAND)
text(s, 1.1, 5.15, 11, 0.9,
     "瑞合瑞德暖通科技集团营销体系数字化平台 —— 覆盖「线索→AI问诊→设计/BIM→报价→合同→施工→验收→IoT 生命周期」经营闭环",
     size=15, color=RGBColor(0x9C, 0xA3, 0xAF), spacing=1.35)
text(s, 1.1, 6.6, 11, 0.4,
     [[("软件厂商 ", {"color": INK3, "size": 12}),
       ("Rysnova / 瑞诺瓦", {"color": WHITE, "size": 12, "bold": True, "font": EN}),
       ("   ·   白标交付客户#1：瑞合瑞德集团   ·   Powered by Rysnova",
        {"color": INK3, "size": 12})]])

# ---- 02 培训议程
s = slide()
page_header(s, "Agenda", "培训议程")
agenda = [
    ("01  产品定位与命名体系", "厘清集团 / 厂商 / 品牌 / 设备的关系"),
    ("02  核心价值闭环", "一条主线看懂平台在做什么"),
    ("03  平台总览架构", "两对外板块 + 增长中枢 + 底座"),
    ("04  板块一 · 品牌管理", "集团官网与设备品牌站群"),
    ("05  板块二 · 经销商赋能", "问诊 / CRM / BIM 三件套 + 四端口"),
    ("06  客户门户与角色权限", "谁登录、看什么、能做什么"),
    ("07  设计流程与六大系统", "双模式设计 + AI 智能能力"),
    ("08  技术架构 · 安全 · 路线图", "选型、多租户隔离、上线规模"),
]
cards(s, agenda, cols=2, top=2.15, bottom=6.75)

# ---- 03 我们解决什么问题
s = slide()
page_header(s, "Why", "我们解决什么问题")
text(s, 0.9, 1.95, 11.5, 0.5,
     "暖通舒适家行业「营销—设计—交付—服务」长期割裂，信息散落、协作低效、客户体验断层。平台把全链路装进一个数字底座。",
     size=14, color=INK2, spacing=1.3)
prob = [
    ("线索散、转化低", "获客渠道分散，缺乏统一沉淀与 AI 归因，跟单靠人肉"),
    ("方案慢、不专业", "现场谈单出方案慢，负荷计算与选型依赖个人经验"),
    ("交付黑盒、客户看不见", "施工/验收进度不透明，客户缺乏可信的项目门户"),
    ("数据孤岛、无生命周期", "签约后设备、质保、服务无系统承接，复购无抓手"),
]
cards(s, prob, cols=2, top=2.75, bottom=6.75)

# ---- 04 产品定位与命名体系
s = slide()
page_header(s, "Identity", "产品定位与命名体系")
text(s, 0.9, 1.95, 11.5, 0.5,
     "一句话：Rysnova（瑞诺瓦）是独立中立的行业软件厂商；平台白标交付给客户，冠客户名并标注 Powered by Rysnova。",
     size=14, color=INK2, spacing=1.3)
idmap = [
    ("瑞合瑞德 / Rhautt Comfort", "暖通科技集团（客户#1）。集团英文/中文表述，不作为软件产品名。"),
    ("Rysnova / 瑞诺瓦", "独立中立的垂直行业软件厂商与经销商赋能体系品牌。"),
    ("Rhautt Nexus / 瑞合数智枢纽", "交付给瑞合瑞德的平台实例名（白标冠客户名）。"),
    ("Rheem · Ruud · Everhot · Lithnova", "设备品牌，进入舒适家方案的产品配置矩阵，各有独立品牌站。"),
]
cards(s, idmap, cols=2, top=2.75, bottom=6.75)

# ---- 05 核心价值闭环
s = slide()
page_header(s, "Closed Loop", "核心价值闭环")
text(s, 0.9, 1.95, 11.5, 0.4, "平台不是松散的功能集合，用一条经营闭环来衡量它的价值：",
     size=14, color=INK2)
steps = ["线索", "AI 问诊", "设计 / BIM", "系统包", "报价", "合同", "施工", "验收", "IoT 生命周期"]
n = len(steps)
left = 0.9; total = SW - 1.8
bw = (total - (n - 1) * 0.18) / n
y = 3.2
for i, st in enumerate(steps):
    x = left + i * (bw + 0.18)
    on = i in (1, 2, 4)   # 高亮 AI 相关环节
    rect(s, x, y, bw, 1.0, fill=(BRAND if on else SURFACE),
         line=(None if on else LINE), line_w=1, rounded=True, radius=0.12)
    text(s, x, y, bw, 1.0, st, size=12.5, color=(WHITE if on else INK),
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < n - 1:
        text(s, x + bw - 0.02, y, 0.2, 1.0, "›", size=18, color=INK3,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 4.6, 11.5, 0.4, "每一步都在同一 RLS 事务内写库并投递 outbox 事件，驱动跨域协同——数据只写一次，全链路自动流转。",
     size=13, color=INK3, spacing=1.3)
low = [
    ("获客侧", "问诊沉淀线索、AI 归因转化"),
    ("经营侧", "设计→报价→合同→施工→验收一站式"),
    ("服务侧", "IoT 生命周期承接，复购与运维闭环"),
]
cards(s, low, cols=3, top=5.2, bottom=6.75)

# ---- 06 平台总览架构
s = slide()
page_header(s, "Architecture", "平台总览架构")
text(s, 0.9, 1.95, 11.5, 0.4, "两个对外产品板块 + 一个对内增长中枢 + 一个共享底座，权限上划分为 6 大域（D0–D5）。",
     size=14, color=INK2)
blocks = [
    ("板块一 · Rhautt 品牌管理", "集团官网 + Everhot / Lithnova 自建站 + Rheem / Ruud 外链 + DAM 物料库", "对外 · D1/D2"),
    ("板块二 · 瑞诺瓦经销商赋能", "问诊 / 舒适家 CRM / 技术支持 BIM 三件套，承载经营闭环", "对外 · D3/D4"),
    ("板块三 · 增长中枢 Nexus Growth", "对内 AI 营销能力域：舆情 / 文案 / GEO / 投放（非对外板块）", "对内 · D5"),
    ("底座 / 总部", "auth · tenant · analytics · governance · notification · workflow · file-artifact", "共享 · D0"),
]
cards(s, blocks, cols=2, top=2.7, bottom=6.75)

# ---- 截图：管理中枢（含局部放大）
shot_slide("Live UI · Console", "界面实拍 · Nexus 管理中枢", "nexus_live.png",
    "apps/nexus-console（:4010）—— 总部/平台管理中枢，侧边栏一屏纵览三大板块（骨架原型·占位数据）。",
    [("三板块导航（左）", "板块一品牌与市场 · 板块二赋能 · 板块三增长"),
     ("板块总览看板", "在管站点 / 物料资产 / 产品条目 / 待发布"),
     ("板块三·增长中枢", "GEO 可见度 / 文案 Copilot / 舆情雷达 / 营销自动化")],
    "/ · 管理中枢",
    zoom={"crop": (0.135, 0.145, 0.02, 0.63), "at": (8.55, 5.75, 3.85),
          "label": "局部放大 · 板块指标卡"})

# ---- 07 板块一 品牌管理
s = slide()
page_header(s, "Board 1", "板块一 · Rhautt 品牌管理")
text(s, 0.9, 1.95, 11.5, 0.4, "对内挂集团/设备品牌。每个品牌站 UI/VI 完全独立，中枢只供给非视觉骨架，不吞并任何独立站点。",
     size=14, color=INK2, spacing=1.3)
b1 = [
    ("集团官网 public-portal", "rhautt.com · 自建，复刻 aosmith.com 内容架构，Ruud 调性"),
    ("Everhot 品牌站", "everhot.com.cn · 自建，复刻 rheem.com 三受众架构"),
    ("Lithnova 储能站", "lithnova.com.cn · 新能源储能设备独立品牌站（自建）"),
    ("Rheem / Ruud", "rheem/ruud.com.cn · 外链占位，不自建"),
    ("品牌控制台 brand-console", "统一管理站群、产品库、DAM 物料、VI token、上新与发布"),
    ("产品目录 product-catalog", "各品牌产品权威源；板块二只读同步选用"),
]
cards(s, b1, cols=3, top=2.7, bottom=6.75)

# ---- 截图：集团官网（实运行 Next.js）
shot_slide("Live UI · Portal", "界面实拍 · 集团官网", "portal_live.png",
    "apps/public-portal（:4005）—— 面向多受众的集团门户，获客与品牌的第一入口。",
    [("顶部导航", "关于我们 / 旗下品牌 / 可持续发展 / 新闻 / 招聘"),
     ("主视觉 Hero", "“创新科技·成就舒适明天” + COP 4.8 能效可视化"),
     ("集团实力数据", "1994 建厂 / 30+ 载历程 / 100+ 品牌基因"),
     ("三大入口", "查找经销商 / 注册产品·保修 / 预约勘测")],
    "rhautt.com")

# ---- 截图：集团官网整页长图
longshot_slide("Live UI · Full Page", "界面实拍 · 集团官网整页长图", "portal_full.png",
    "完整首页滚动全貌（Playwright fullPage 抓取），展示内容编排与信息层次。",
    [("首屏 Hero", "品牌主张 + 能效数据可视化"),
     ("业务与品牌", "四大领域与旗下设备品牌展示"),
     ("可持续 / 新闻", "ESG 叙事与集团动态"),
     ("页脚导航", "多入口与联系方式")],
    "rhautt.com")

# ---- 08 板块二 总览 + 四端口
s = slide()
page_header(s, "Board 2", "板块二 · 瑞诺瓦经销商赋能（总览）")
text(s, 0.9, 1.95, 11.5, 0.4, "对外以「瑞诺瓦 / Rysnova」中立行业软件形态呈现。三件套 + 四个前端端口应用。",
     size=14, color=INK2, spacing=1.3)
text(s, 0.9, 2.55, 11.5, 0.3, "三件套", size=13, color=BRAND, bold=True, font=EN)
suites = [
    ("① AI 问诊", "C 端业主智能问诊获客", "diagnosis"),
    ("② 舒适家 CRM", "经销商经营闭环主战场", "crm/quote/delivery"),
    ("③ 技术支持 BIM", "设计师 / BIM 深化选型", "rysnova-bim/design"),
]
cards(s, suites, cols=3, top=2.9, bottom=4.55)
text(s, 0.9, 4.75, 11.5, 0.3, "四个前端端口应用", size=13, color=BRAND, bold=True, font=EN)
ports = [
    ("集团门户 :4005", "对外获客入口"),
    ("经销商工作台 :4000", "销售/店长·登录中心"),
    ("设计师工作台 :4003", "设计/技术支持"),
    ("客户门户 :4002", "签约客户看进度"),
]
cards(s, ports, cols=4, top=5.1, bottom=6.75)

# ---- 截图：统一登录入口
shot_slide("Live UI · Login", "界面实拍 · 统一登录入口", "dealer_live.png",
    "apps/dealer-workbench（:4000）—— 登录中心：登录一次，按角色进入所有应用。",
    [("品牌面板", "Rhautt Comfort · Powered by Rysnova AI"),
     ("能力总览", "品牌厂家功能组 / 赋能三件套 / 平台运营 / 客户入口"),
     ("统一登录", "账号密码登录，SSO cookie 贯通各工作台"),
     ("多入口跳转", "返回集团门户 / 经营控制台登录")],
    "/ · 登录中心")

# ---- 09 件套① AI问诊
s = slide()
page_header(s, "Suite 1", "件套① · 瑞诺瓦 AI 问诊")
text(s, 0.9, 1.95, 11.5, 0.5, "面向 C 端业主的智能问诊入口（apps/consumer-diagnosis），把痛点诊断转化为高质量线索。",
     size=14, color=INK2, spacing=1.3)
d = [
    ("痛点智能诊断", "对话式采集户型、需求与痛点，AI 生成初步舒适家诊断报告"),
    ("方案初判与引导", "输出六大系统适配建议，引导预约到店/上门与深度设计"),
    ("线索沉淀归因", "诊断即留资，自动进入 CRM 线索池并归因到渠道/门店"),
    ("多入口嵌入", "可嵌入集团官网、各品牌站，保持中立工具形态对外"),
]
cards(s, d, cols=2, top=2.85, bottom=6.75)

# ---- 截图：AI 问诊
shot_slide("Live UI · Suite 1", "界面实拍 · 瑞诺瓦 AI 问诊", "diagnosis_live.png",
    "apps/consumer-diagnosis（:4001）—— 对外以中立“瑞诺瓦 / RYSNOVA”形态，把痛点转化为线索。",
    [("品牌导航", "瑞诺瓦 RYSNOVA 中立行业形态对外"),
     ("主视觉", "AI-POWERED HOME COMFORT，3 分钟生成专属方案"),
     ("免费 AI 问诊 CTA", "一键进入对话式问诊流程"),
     ("六大系统分区", "热水 / 空气与净水等系统方案入口")],
    "rysnova · AI 问诊")

# ---- 10 件套② 舒适家CRM
s = slide()
page_header(s, "Suite 2", "件套② · 瑞诺瓦舒适家 CRM")
text(s, 0.9, 1.95, 11.5, 0.5, "经销商 / 门店 / 销售 / 已签约客户的经营主战场，承载报价→合同→施工→验收→IoT 生命周期。",
     size=14, color=INK2, spacing=1.3)
c = [
    ("线索 → 商机 → 签单", "全漏斗跟单，商机、报价、合同状态一屏可见"),
    ("智能报价系统", "多促销配置、实时算价、材料清单与报价单导出"),
    ("交付与验收", "施工过程、节点验收记录，产物经 file-artifact 存储"),
    ("生命周期 / IoT 移交", "签约客户、家庭、设备、质保、服务计划完整承接"),
]
cards(s, c, cols=2, top=2.85, bottom=6.75)

# ---- 截图：经营工作台
shot_slide("Live UI · Suite 2", "界面实拍 · 经营工作台（CRM）", "console.png",
    "多租户业务工作台，经销商/门店/总部经营全局一屏可见。",
    [("经营总览导航", "报价测算 · 产品 · 促销 · 施工闭环 · 客户经营"),
     ("销售漏斗 / 经销商表现", "全局经营看板与排名"),
     ("阶段金额分布", "商机金额分布与转化分析"),
     ("总部闭环预警", "各阶段项目分布与卡点预警")],
    "/business-console")

# ---- 11 客户门户（新增，回应用户）
s = slide()
page_header(s, "Customer Portal", "客户门户 · 签约客户端口")
text(s, 0.9, 1.95, 11.5, 0.6,
     "apps/customer-portal（:4002）—— 隶属板块二 CRM 的 lifecycle/delivery。签约客户登录后，"
     "可透明查看与自己项目相关的全部信息，让交付不再是黑盒。",
     size=14, color=INK2, spacing=1.3)
cp = [
    ("沟通记录", "与经销商/设计师的沟通往来集中留痕，随时可查"),
    ("报价与合同", "查看专属报价单、方案与合同，明明白白消费"),
    ("资料与图纸", "系统图、施工图、设备清单等项目资料在线查阅"),
    ("服务与进度", "施工节点、验收状态、质保与服务计划实时可视化"),
]
cards(s, cp, cols=2, top=3.0, bottom=6.75)

# ---- 截图：客户门户
shot_slide("Live UI · Customer", "界面实拍 · 客户门户", "customer_live.png",
    "apps/customer-portal（:4002）—— 签约客户登录端口，让交付全程透明可见。",
    [("品牌标识", "瑞诺瓦舒适家 · 客户服务门户"),
     ("分享码 / 手机号登录", "无需密码，轻量化进入专属方案"),
     ("查看专属方案", "报价、方案、施工与验收进度")],
    "/customer")

# ---- 12 件套③ 技术支持BIM
s = slide()
page_header(s, "Suite 3", "件套③ · 瑞诺瓦技术支持 BIM")
text(s, 0.9, 1.95, 11.5, 0.5, "面向设计师 / 技术支持（apps/rysnova-bim-workbench + designer-workbench），是平台的专业护城河。",
     size=14, color=INK2, spacing=1.3)
b = [
    ("BIM 深化设计", "承接签单方案，进行专业 BIM 深化与系统图/施工图输出"),
    ("Revit / CAD 集成", "对接 Revit 插件与 CAD，2D 设计选型到三维深化打通"),
    ("产物与签收", "深化产物经 file-artifact/对象存储管理，交付可签收"),
    ("与设计放行联动", "design/ai-design 精算放行，交付/生命周期串联闭环"),
]
cards(s, b, cols=2, top=2.85, bottom=6.75)

# ---- 截图：Rysnova BIM 工作台
shot_slide("Live UI · Suite 3", "界面实拍 · Rysnova BIM 工作台", "bim.png",
    "面向设计师/技术支持的专业 BIM 深化工作台，是平台的专业护城河。",
    [("项目与建筑参数", "面积 / 楼层 / 层高 / 系统选择"),
     ("3D·平面·立面视图", "多视角专业深化设计"),
     ("设计数据面板", "负荷计算 / 设备清单 / 管道统计 / 规范合规"),
     ("产物导出", "一键导出 IFC / DWG / 3D 模型")],
    "/rysnova-bim-designer")

# ---- 13 双模式设计流程
s = slide()
page_header(s, "Design Flow", "双模式设计流程")
text(s, 0.9, 1.95, 11.5, 0.4, "同一套引擎，两种节奏：现场快速谈单，回来精细化施工级设计。",
     size=14, color=INK2)
# 两大模式卡
rect(s, 0.9, 2.6, 5.55, 3.9, fill=SURFACE, line=LINE, rounded=True, radius=0.04)
rect(s, 0.9, 2.6, 5.55, 0.75, fill=BRAND, rounded=True, radius=0.06)
text(s, 1.2, 2.6, 5.0, 0.75, "快速估算模式", size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.2, 3.55, 5.0, 2.8, [
    [("3–5 分钟出方案 · 适合现场谈单", {"size": 13, "color": INK2, "bold": True})],
    [("· 基于户型与需求的 AI 推荐", {"size": 13, "color": INK2})],
    [("· 简易 3D 展示", {"size": 13, "color": INK2})],
    [("· 预估报价快速输出", {"size": 13, "color": INK2})],
], spacing=1.5, space_after=8)
rect(s, 6.9, 2.6, 5.55, 3.9, fill=SURFACE, line=LINE, rounded=True, radius=0.04)
rect(s, 6.9, 2.6, 5.55, 0.75, fill=INK, rounded=True, radius=0.06)
text(s, 7.2, 2.6, 5.0, 0.75, "精细化设计模式", size=18, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
text(s, 7.2, 3.55, 5.0, 2.8, [
    [("30 分钟出施工级设计", {"size": 13, "color": INK2, "bold": True})],
    [("· 户型绘制：手绘 / CAD 导入 / 模板", {"size": 13, "color": INK2})],
    [("· 专业负荷计算（符合国家标准）", {"size": 13, "color": INK2})],
    [("· 智能设备选型 + 3D 布局", {"size": 13, "color": INK2})],
    [("· 完整材料清单 + 精准报价", {"size": 13, "color": INK2})],
], spacing=1.5, space_after=8)

# ---- 截图：设计师工作台
shot_slide("Live UI · Designer", "界面实拍 · 设计师工作台", "designer.png",
    "全屋舒适系统 2D 交互设计与选型，边设计边算价。",
    [("工具栏", "设备 / 智能 / 布局等绘制工具"),
     ("画布", "平面交互体系：房型 / 设备 / 管线 / 材料"),
     ("对象属性", "选中设备的参数编辑"),
     ("实时材料清单与报价", "自动估价，一键生成客户报价")],
    "/designer")

# ---- 14 六大系统 + AI能力
s = slide()
page_header(s, "Systems & AI", "六大系统 + AI 智能能力")
text(s, 0.9, 1.95, 11.5, 0.35, "六大舒适家系统", size=13, color=BRAND, bold=True, font=EN)
sys6 = [
    ("五恒系统", "恒温恒湿恒氧"), ("净水系统", "全屋净水解决方案"),
    ("采暖系统", "节能采暖方案"), ("热水系统", "中央热水供应"),
    ("新风系统", "全热交换新风"), ("除湿系统", "智能除湿控制"),
]
cards(s, sys6, cols=3, top=2.3, bottom=4.55)
text(s, 0.9, 4.75, 11.5, 0.35, "AI 智能能力", size=13, color=BRAND, bold=True, font=EN)
ai = [
    ("智能负荷计算", "按国标自动算热/冷负荷"),
    ("设备自动选型", "匹配品牌产品矩阵"),
    ("AI 方案推荐 / 3D 布局", "自动布局 + 碰撞检测优化"),
]
cards(s, ai, cols=3, top=5.1, bottom=6.75)

# ---- 15 板块三 增长中枢（新增）
s = slide()
page_header(s, "Board 3", "板块三 · 增长中枢 Nexus Growth")
text(s, 0.9, 1.95, 11.5, 0.6,
     "对内 AI 营销能力域（/api/v2/growth）——「用魔法打败魔法」，把 AI 装进合规、透明、可审计的营销工作流。"
     "只采公开数据，AI 产出默认待人工核准。",
     size=14, color=INK2, spacing=1.3)
g = [
    ("E1 舆情监测 Sentiment", "全网公开源实时监测 + AI 情感/意图分级 + 危机预警", "E1"),
    ("E2 文案策划 Copilot", "多平台风格文案矩阵一键生成 + 品牌护栏 + 人工核准", "E2"),
    ("E3 GEO 分析 Analyzer", "生成式引擎优化：让 AI 搜索优先引用我方权威内容", "E3"),
    ("E4 营销自动化 Ops", "线索归因 + ROI 闭环 + 素材 A/B + 预算再分配建议", "E4"),
]
cards(s, g, cols=2, top=3.05, bottom=6.75)

# ---- 16 技术架构
s = slide()
page_header(s, "Tech Stack", "技术架构")
text(s, 0.9, 1.95, 11.5, 0.4, "终态锁定：NestJS + Fastify + PostgreSQL（RLS 多租户）；后端唯一真相源 /api/v2/*。",
     size=14, color=INK2, spacing=1.3)
tech = [
    ("前端", "React 18 + TypeScript · Vite · Tailwind · Three.js/R3F · Zustand · React Query；Next.js 多应用 monorepo"),
    ("后端（终态）", "NestJS + Fastify（DDD 模块化单体，services/api）· JWT（HttpOnly cookie）"),
    ("数据层", "PostgreSQL + TypeORM（RLS 行级隔离）· MongoDB 文档库 · Redis 缓存/会话"),
    ("编排 / 事件", "Temporal + Outbox 跨域事件（写业务与写 outbox 同事务原子）"),
    ("认证 / 网关", "cookie SSO（nx_token, shared-auth）· Express 反向代理 → NestJS"),
    ("工程 / 部署", "ESLint + Prettier · Docker · Nx · CI 门禁强制架构与品牌规范"),
]
cards(s, tech, cols=2, top=2.7, bottom=6.75)

# ---- 17 角色与权限
s = slide()
page_header(s, "Roles", "角色与权限")
text(s, 0.9, 1.95, 11.5, 0.4, "多角色协同，权限按 tenant / dealer / store / role 分层，配合 6 大权限域（D0–D5）。",
     size=14, color=INK2, spacing=1.3)
roles = [
    ("总部 / 管理员", "跨租户汇总分析、治理、品牌与产品配置（business-console）"),
    ("经销商 / 店长", "门店经营、团队与线索分配、报价与合同审批"),
    ("销售", "跟单、报价、签约，客户信息管理"),
    ("设计师 / 技术支持", "方案设计、BIM 深化、施工图与选型"),
    ("市场 / 增长", "板块三增长中枢：舆情、文案、GEO、投放（hq_marketing）"),
    ("签约客户 / 业主", "客户门户查看沟通、报价、资料与服务进度"),
]
cards(s, roles, cols=2, top=2.7, bottom=6.75)

# ---- 18 数据安全与多租户
s = slide()
page_header(s, "Security", "数据安全与多租户隔离")
text(s, 0.9, 1.95, 11.5, 0.5,
     "面向 500+ 经销商租户，安全与合规是一等公民。板块级物理隔离 + 租户级 RLS 双层叠加。",
     size=14, color=INK2, spacing=1.3)
sec = [
    ("行级安全 RLS", "所有租户数据经 Postgres RLS + withRlsTransaction 隔离"),
    ("双层隔离", "物理分库守板块边界，RLS 守租户边界，二者叠加不互替"),
    ("PII 列级加密", "静态加密 + PII 列级加密 + 审计留痕，符合中国合规（PIPL）"),
    ("白标交付", "私有部署冠客户名，分析层脱敏，交付实例数据互不可见"),
]
cards(s, sec, cols=2, top=2.85, bottom=6.75)

# ---- 19 上线规模与路线图
s = slide()
page_header(s, "Scale & Roadmap", "上线规模与演进路线")
text(s, 0.9, 1.95, 11.5, 0.35, "目标规模", size=13, color=BRAND, bold=True, font=EN)
scale = [
    ("500+", "经销商租户并发"),
    ("2000+", "设计 / 销售人员"),
    ("10 万+", "用户 / 客户档案"),
]
cards(s, scale, cols=3, top=2.3, bottom=4.0)
text(s, 0.9, 4.2, 11.5, 0.35, "迁移演进路线（不大爆炸，按域推进）", size=13, color=BRAND, bold=True, font=EN)
road = [
    ("Step 1 · auth", "认证域优先迁至 NestJS/PG"),
    ("Step 2 · tenant", "多租户与 RLS 落地"),
    ("Step 3 · crm", "经营闭环主战场迁移"),
    ("Step 4 · quote", "报价与后续域逐步下线遗留"),
]
cards(s, road, cols=4, top=4.55, bottom=6.75)

# ==== 瑞诺瓦 ISV 产品专题 ==========================================
# ---- ISV-A 产品定位
s = slide()
page_header(s, "Rysnova ISV", "瑞诺瓦 · 舒适家行业软件产品")
text(s, 0.9, 1.95, 11.5, 0.6,
     "瑞诺瓦（Rysnova）是面向暖通舒适家行业的独立中立 ISV：对客以「经销商联合主体 + 中立行业工具」形态呈现，"
     "平台白标交付、冠客户名 · Powered by Rysnova。",
     size=14, color=INK2, spacing=1.3)
isv = [
    ("产品形态", "多租户 SaaS + 私有化白标双轨交付，一套内核服务全行业", "SaaS / 白标"),
    ("目标客户", "暖通舒适家集成商、经销商、门店与设计 / 技术支持团队", "B 端经销商"),
    ("一体化范围", "AI 问诊获客 · 舒适家 CRM 经营 · 技术支持 BIM，覆盖全链路", "三件套"),
    ("交付主张", "30 分钟出施工级方案 · 自动精算报价 · 整包质保 · 交付透明", "价值差"),
]
cards(s, isv, cols=2, top=3.0, bottom=6.75)

# ---- ISV-B 五大结构性特点 / 护城河
s = slide()
page_header(s, "Product Moat", "结构性产品特点与护城河")
text(s, 0.9, 1.95, 11.5, 0.4,
     "窗口期优势易被复制，真正难复制的是「标准级计算 + 物料同源 + 行业情报网」三项结构性壁垒。",
     size=14, color=INK2, spacing=1.3)
moat = [
    ("标准级计算内核", "工程内核 + 国标 / ASHRAE 精度基线，计算可回归、可追责"),
    ("物料同源", "UI / VI / SI 与市场物料从设计系统同源生成，竞品无此能力"),
    ("产品 / 价格情报网", "跨租户产品主数据 + 调用量沉淀为行业级价格与选型情报"),
    ("多品牌中枢", "Rheem / Ruud / Everhot 等多设备品牌统一中枢（竞品多为单品牌）"),
    ("全链路闭环", "问诊→报价→合同→设计→BIM→施工→验收→IoT 一条数据链贯通"),
    ("单一真相源", "design ↔ BIM 版本冻结 / 审图签章 / 回滚，交付可追责不分叉"),
]
cards(s, moat, cols=2, top=2.7, bottom=6.75)

# ---- ISV-C 核心功能价值（三件套）
s = slide()
page_header(s, "Functional Value", "核心功能价值 · 赋能三件套")
suite_val = [
    ("① AI 问诊获客", "对话式舒适体检，痛点→高质量线索，自动归因沉淀进 CRM", "C 端获客"),
    ("② 舒适家 CRM", "线索→商机→报价→合同→施工→验收→IoT 全漏斗经营闭环", "经营中台"),
    ("③ 技术支持 BIM", "负荷精算 / 选型 / 3D 深化 / 碰撞检测 / 管线综合 / IFC·DWG", "专业护城河"),
]
cards(s, suite_val, cols=3, top=2.25, bottom=4.4)
text(s, 0.9, 4.6, 11.5, 0.35, "对经销商的三重价值", size=13, color=BRAND, bold=True, font=EN)
tri = [
    ("提效", "30 分钟出施工级设计，报价自动精算，告别 Excel 拍脑袋"),
    ("提质", "国标合规计算 + 整包质保，交付全程透明，返工与纠纷下降"),
    ("提转化", "AI 问诊获客 + 全链路归因，线索转化率与客单价可衡量提升"),
]
cards(s, tri, cols=3, top=5.05, bottom=6.75)

# ---- ISV-D 行业价值评估：痛点 → 解法
s = slide()
page_header(s, "Industry Value", "行业价值评估 · 痛点 → 瑞诺瓦解法")
text(s, 1.15, 2.2, 4.7, 0.3, "行业痛点", size=12, color=INK3, bold=True, font=EN)
text(s, 6.5, 2.2, 5.7, 0.3, "瑞诺瓦解法", size=12, color=BRAND, bold=True, font=EN)
pains = [
    ("报价靠经验、一物多价", "标准级精算 + 价格快照锁定，杜绝一物两价"),
    ("设计慢、重度依赖资深工程师", "双模式设计，30 分钟出施工级方案，能力下沉"),
    ("交付黑盒、客户缺信任", "客户门户项目护照，施工 / 验收全程透明可视"),
    ("多品牌多系统割裂", "多品牌中枢 + 六大舒适家系统统一配置选型"),
    ("物料散乱、品牌调性不统一", "物料同源，UI / VI / SI 与营销物料一键生成"),
    ("无数据沉淀、决策拍脑袋", "全链路数据 + 跨租户行业情报，量化经营决策"),
]
ry = 2.6
for p, sol in pains:
    rect(s, 0.9, ry, 11.5, 0.58, fill=SURFACE, line=LINE, rounded=True, radius=0.06)
    text(s, 1.15, ry, 4.6, 0.58, p, size=13, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.75, ry, 0.6, 0.58, "→", size=17, color=BRAND, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 6.5, ry, 5.75, 0.58, sol, size=12.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.68

# ---- ISV-E 市场潜力与增长飞轮
s = slide()
page_header(s, "Market Potential", "市场潜力与增长飞轮")
text(s, 0.9, 1.95, 11.5, 0.55,
     "舒适家 / 暖通集成是万亿级家装升级赛道，行业数字化渗透率低、工具化空白大，ISV 复制空间广阔。",
     size=14, color=INK2, spacing=1.3)
fly = ["更多经销商入驻", "更多项目与数据", "精算 / 情报更强", "价值差更大"]
nf = len(fly)
lf = 0.9; totf = SW - 1.8
bwf = (totf - (nf - 1) * 0.55) / nf
yf = 2.95
for i, st in enumerate(fly):
    x = lf + i * (bwf + 0.55)
    rect(s, x, yf, bwf, 0.95, fill=(BRAND if i == 0 else SURFACE),
         line=(None if i == 0 else LINE), line_w=1, rounded=True, radius=0.12)
    text(s, x, yf, bwf, 0.95, st, size=13, color=(WHITE if i == 0 else INK),
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < nf - 1:
        text(s, x + bwf + 0.02, yf, 0.5, 0.95, "↻", size=18, color=BRAND,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 4.05, 11.5, 0.35, "网络效应自我强化：租户越多 → 数据 / 情报越强 → 价值差越大 → 迁移动力越强，形成正循环。",
     size=12.5, color=INK3, spacing=1.3)
pot = [
    ("网络效应", "租户越多，产品 / 价格情报越强，后来者越难复制"),
    ("数据反哺", "IoT 运行数据反向校准精算模型（生命周期 → 设计）"),
    ("可复制 ISV", "白标 + 多租户，快速复制到更多集团 / 经销商 / 设备品牌"),
    ("增量变现", "订阅 + 交付增值 + 情报 / 供应链金融等衍生服务"),
]
cards(s, pot, cols=2, top=4.65, bottom=6.75)

# ---- 20 结尾
s = slide(INK)
rect(s, 0, 0, SW, SH, fill=INK)
rect(s, 0, 0, 0.22, SH, fill=BRAND)
text(s, 1.1, 2.5, 11, 1.2, [
    [("一个平台，跑通舒适家全链路经营闭环", {"color": WHITE, "size": 34, "bold": True})],
], spacing=1.1)
rect(s, 1.1, 3.7, 0.9, 0.07, fill=BRAND)
text(s, 1.1, 3.95, 11, 1.0,
     "线索 → AI 问诊 → 设计/BIM → 系统包 → 报价 → 合同 → 施工 → 验收 → IoT 生命周期",
     size=15, color=RGBColor(0x9C, 0xA3, 0xAF), spacing=1.4)
text(s, 1.1, 5.3, 11, 0.4, "谢谢观看 · Thank You", size=20, color=WHITE, bold=True, font=EN)
text(s, 1.1, 6.6, 11, 0.4,
     "Rhautt Nexus / 瑞合数智枢纽   ·   Powered by Rysnova",
     size=12, color=INK3)

prs.save("Rhautt-Nexus-培训.pptx")
print("OK ->", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
